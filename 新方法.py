from docx import Document
from docx.enum.table import WD_ALIGN_VERTICAL
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls as ptspace
from pptx.oxml.xmlchemy import  OxmlElement as pptXML
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT,MSO_AUTO_SIZE,MSO_VERTICAL_ANCHOR
from pptx.util import  Pt
import html
from io import BytesIO
from modules.get_pdf_info import get_pdf_page_info
from modules.oprate_word import find_all_elements
import os
from docx.oxml.ns import qn as wdqn
import xmltodict
import dicttoxml
import json

from docx.oxml import parse_xml as wdparse_xml
from docx.oxml.ns import nsdecls  as wdspace 
from lxml import etree


doc=Document("trainning.docx")
prs = Presentation("test.pptx")


oMaths= find_all_elements(doc._element, wdqn('m:oMath'))
for oMath in oMaths:
  runele=wdparse_xml(f'''<w:r {wdspace("w")}>{etree.tostring(oMath)}</w:r>''')
  if oMath.getparent().tag==wdqn("m:oMathPara"):
    oMath.getparent().getparent().replace(oMath.getparent(),runele)
  else:
    oMath.getparent().replace(oMath,runele)


word_para=doc.paragraphs[0]
ppt_para  = prs.slides[0].shapes[0].text_frame.paragraphs[0] 


word_dict =xmltodict.parse(word_para._element.xml)
ppt_dict=xmltodict.parse(ppt_para._element.xml)


open('a.json','w',encoding='utf-8').write(json.dumps(word_dict,ensure_ascii=False,indent=4))


open('b.json','w',encoding='utf-8').write(json.dumps(ppt_dict,ensure_ascii=False,indent=4))



exit()


def set_para_format(word_para,ppt_para):
  word_dict =xmltodict.parse(word_para._element.xml)
  ppt_dict=xmltodict.parse(ppt_para._element.xml)

  #设置段落对齐方式
  w_jc=word_dict["w:p"]["w:pPr"].get('w:jc')
  if w_jc:
      ppt_dict["a:p"]["a:pPr"]= {'@algn': {'left':'l','center':'ctr','right':'r','both':'just','distribute':'dist'}[w_jc['@w:val']]}

  w_ind=word_dict["w:p"]["w:pPr"].get('w:ind')
  if w_ind:
    #设置文本之前，文本之后的缩进
    if w_ind.get('@w:left'):
      ppt_dict["a:p"]["a:pPr"]['@marL']=f'{int(w_ind.get('@w:left'))*635}'
    if w_ind.get('@w:right'):
      ppt_dict["a:p"]["a:pPr"]['@marR']=f'{int(w_ind.get('@w:right'))*635}'
    #设置段落首行缩进
    if w_ind.get('@w:firstLine'):
      ppt_dict["a:p"]["a:pPr"]['@indent']=f'{int(w_ind.get('@w:firstLine'))*635}'
    #设置段落悬挂缩进
    if w_ind.get('@w:hanging'):
      ppt_dict["a:p"]["a:pPr"]['@indent']=f'-{int(w_ind.get('@w:hanging'))*635}'
      
  w_spacing=word_dict["w:p"]["w:pPr"].get("w:spacing")
  if w_spacing:
    #设置段前、段后间距
    if w_spacing.get('@w:beforeLines'):
      ppt_dict["a:p"]["a:pPr"]['a:spcBef']={"a:spcPts": {"@val": f'{int(w_spacing.get("@w:beforeLines"))*15.6}'}}
    if  w_spacing.get('@w:afterLines'):
      ppt_dict["a:p"]["a:pPr"]['a:spcAft']={"a:spcPts": {"@val": f'{int(w_spacing.get("@w:afterLines"))*15.6}'}}
    #设置行距
    if w_spacing.get('"@w:line'):
      ppt_dict["a:p"]["a:pPr"]['a:lnSpc']={"a:spcPct": {"@val": f'{int(w_spacing.get("@w:line"))*1250/3}'}}

  #设置制表位    
  w_tabs=word_dict["w:p"]["w:pPr"].get("w:tabs")
  if w_tabs:
    ppt_dict["a:p"]["a:pPr"]['a:tabLst']={"a:tab": [{"@pos": f'{int(i["@w:pos"])*635}',"@algn": "l"} for i in w_tabs['w:tab']]}


prs.save('test.pptx') 

