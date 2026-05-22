from docx import Document
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import  Pt
from io import BytesIO
from modules.get_pdf_info import get_info
from modules.oprate_word import *
import os
def docx2pptx(docfile,bookname,pptfolder):

  doc = Document(docx=docfile)
  prs = Presentation("modules\母版.pptx")
  log=open(docfile.replace(".docx",".log"),'a',encoding='utf-8')
  for shape in prs.slide_master.shapes:
    if shape.has_text_frame:
        shape.text_frame.paragraphs[0].runs[0].text=bookname

  paragraphs=doc.paragraphs
  tables=doc.tables
  page_infos=get_info(docfile)

  filepath=[]
  new=False

  for num,page_info in page_infos.items(): 
    slide = prs.slides.add_slide(prs.slide_layouts[0]) #pdf的一页对应一个幻灯片
    for content in page_info['contents']:
      txBox = slide.shapes.add_textbox(Pt(28), Pt(content[1]),Pt(664), Pt(content[3])) #默认版心宽，常量
      text_frame= txBox.text_frame
      text_frame.auto_size,text_frame.word_wrap= MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT, True
      (text_frame.margin_bottom,text_frame.margin_top,text_frame.margin_left,
    text_frame.margin_right)=(0,0,0,0)  #文字和文本框间距

      paragraph_text=''
      while not  paragraph_text.startswith('^'):
        paragraph = paragraphs.pop(0)
        paragraph_text= paragraph.text

      if paragraph.style.name =='Heading 1':
        title1,title2,title3=paragraph.text,'',''
        new=True
      elif paragraph.style.name =='Heading 2':
        title2,title3=paragraph.text,''
        new=True
      elif paragraph.style.name =='Heading 3':
        title3=paragraph.text
        new=True
      else:
        if new:
          filepath=[]
          for a in [title1,title2,title3]:
            if a:
              filepath.append(a.replace('^','').replace('*','').replace('?','').replace(':','').replace('"','').replace('<','').replace('>','').replace('|','').replace('/','').replace('\\','').replace('\t',' ').replace('\u3000',' '))
          new=False

      if filepath and new:
        
        try:
          os.makedirs(os.path.join(pptfolder,'\\'.join(filepath[:-1]))) 
        except :
            pass
        prs.save(os.path.join(pptfolder,'\\'.join(filepath)+'.pptx'))
        
        prs = Presentation("modules\母版.pptx")
        for shape in prs.slide_master.shapes:
          if shape.has_text_frame:
              shape.text_frame.paragraphs[0].runs[0].text=bookname

        slide = prs.slides.add_slide(prs.slide_layouts[0]) #pdf的一页对应一个幻灯片
        txBox = slide.shapes.add_textbox(Pt(28), Pt(content[1]),Pt(664), Pt(content[3])) #默认版心宽，常量
        text_frame= txBox.text_frame
        text_frame.auto_size,text_frame.word_wrap= MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT, True
        (text_frame.margin_bottom,text_frame.margin_top,text_frame.margin_left,
        text_frame.margin_right)=(0,0,0,0)  #文字和文本框间距
        filepath=[]

      p  = text_frame.paragraphs[0] 
      set_paragraph(p,paragraph,content[4],1.2) # 设置段落样式
      for ele in paragraph._element:
        a=addelement(p,ele)

    for content in page_info['tables']:
      table=tables.pop(0)
      newtable = slide.shapes.add_table(len(table.rows),len(table.columns),Pt(content[0]), Pt(content[1]),Pt(content[2]), Pt(content[3])).table
      set_table(newtable,table,content[4],num,log)

    for content in page_info['imgs']:
      slide.shapes.add_picture(BytesIO(content[4]),Pt(content[0]), Pt(content[1]),Pt(content[2]), Pt(content[3]))

  prs.save(os.path.join(pptfolder,'\\'.join(filepath)+'.pptx'))
  log.close()