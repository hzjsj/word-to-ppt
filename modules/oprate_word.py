import re
import html
from io import BytesIO
from lxml import etree

from docx import Document
from docx.oxml.ns import qn as wdqn
from docx.oxml.ns import nsdecls  as wdspace  
from docx.oxml import parse_xml as wdparse_xml
from docx.shared import  Pt as wdPt
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.table import WD_ALIGN_VERTICAL

from modules.mathml2latexmaster.mathml2latex import omath_convert

from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls as ptspace
from pptx.oxml.xmlchemy import  OxmlElement as pptXML
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT,MSO_AUTO_SIZE,MSO_VERTICAL_ANCHOR
from pptx.util import  Pt

def find_all_elements(element, target_tag):
    '''按照节点名查找所有下级元素'''
    result = []
    # 查找当前元素的所有子元素
    children = element.findall('./*')
    for child in children:
        if child.tag == target_tag:
            result.append(child)
        # 递归查找子元素的子元素
        else:
          result.extend(find_all_elements(child, target_tag))
    return result
def modify_json_node(data, target_key, new_value):
    if isinstance(data, dict):
        for key, value in data.items():
            if key == target_key:
                data[key] = new_value
            else:
                modify_json_node(value, target_key, new_value)
    elif isinstance(data, list):
        for item in data:
            modify_json_node(item, target_key, new_value)
    return data
#向ppt的一个段落中插入一个word元素
def addelement(p,ele):
  if  ele.tag==wdqn("w:r") and ele.text!="^" and ele.text:
    new=parse_xml(f'''<a:r {ptspace("a")}><a:rPr lang="en-US" altLang="zh-CN" sz="{int(ele.rPr.sz.get(wdqn("w:val")))*50}"><a:latin typeface="{ele.rPr.rFonts.get(wdqn("w:ascii"))}" /> <a:ea typeface="{ele.rPr.rFonts.get(wdqn("w:eastAsia"))}" /></a:rPr><a:t>{html.escape(ele.text)}</a:t></a:r> ''')
    new.rPr.set("b","1") if  ele.rPr.xpath("w:b") else None #如果找到就加粗
    new.rPr.set("i","1") if ele.rPr.xpath("w:i") else None#如果找到就倾斜
    if  ele.rPr.xpath("w:color"): #如果设置了颜色      
      new.rPr.append(parse_xml(f'''<a:solidFill {ptspace("a")}>
        <a:srgbClr val="{ele.rPr.color.get(wdqn("w:val"))}"/>
        </a:solidFill>'''))
    wu_pu={"single":"sng", "double":"dbl","dottedHeavy":"dottedHeavy","wave":"wavy"}
    if ele.rPr.xpath("w:u"): #如果找到就加下划线
      try:
        new.rPr.set("u",wu_pu[ele.rPr.xpath("w:u")[0].get(wdqn("w:val"))])
      except:
        pass
      # ele.rPr.xpath("w:u")[0].get(wdqn("w:color"))  #word下划线颜色，暂默认
      new.rPr.append(parse_xml(f'''<a:uFill {ptspace("a")}><a:solidFill> <a:srgbClr val="000000"/></a:solidFill> </a:uFill>''')) #{ele.rPr.xpath("w:u")[0].get(wdqn("w:color"))}
    
    if ele.rPr.xpath("w:vertAlign"): #如果找到就设置上标，下标
      if ele.rPr.xpath("w:vertAlign")[0].get(wdqn("w:val"))=="superscript":
        new.rPr.set("baseline","30000")
      elif ele.rPr.xpath("w:vertAlign")[0].get(wdqn("w:val"))=="subscript":
        new.rPr.set("baseline","-25000")
    p._element.append(new)
  
  elif  ele.tag==wdqn("m:oMathPara"):
    ele=ele.xpath("./*[1]")[0]
    addelement(p,ele)

  elif  ele.tag==wdqn("m:oMath"):
    #对每个样式重新替换
    for rPr in find_all_elements(ele, wdqn('w:rPr')):
      newrPr=parse_xml(f'''<a:rPr {ptspace('a')} lang="en-US" altLang="zh-CN" sz="{int(rPr.sz.get(wdqn("w:val")))*50}"><a:latin typeface="Cambria Math"/><a:ea typeface="宋体"/></a:rPr> ''' )
      
      if  rPr.xpath("w:color"): #如果设置了颜色      
        newrPr.append(parse_xml(f'''<a:solidFill {ptspace("a")}> <a:srgbClr val="{rPr.color.get(wdqn("w:val"))}"/></a:solidFill>'''))

      newrPr.set("i","1") if rPr.xpath("w:i") else None#如果找到就倾斜
      rPr.getparent().replace(rPr,newrPr)

      # rPr.set("b","1") if  rPr.xpath("w:b") else None #如果找到就加粗
        # wu_pu={"single":"sng", "double":"dbl","dottedHeavy":"dottedHeavy","wave":"wavy"}
        # if ele.xpath("w:u"): #如果找到就加下划线
        #   try:
        #     new.set("u",wu_pu[ele.xpath("w:u")[0].get(wdqn("w:val"))])
        #   except:
        #     pass
    child = parse_xml('''<a14:m xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main"></a14:m>''')
    child.append(ele)
    p._element.append(child)
    return True
def set_paragraph(p,paragraph,indent,line_spacing):
  #段落对齐方式
  try:
    paragraph_aliment=paragraph._element.pPr.jc.get(wdqn("w:val"))#原始段落对齐方式的值
  except AttributeError:
    paragraph_aliment='left'
  p_aliment={'left':'l','center':'ctr','right':'r','both':'just','distribute':'dist'}[paragraph_aliment] #新段落对齐方式的值
  p.alignment=PP_PARAGRAPH_ALIGNMENT.CENTER #新建个标签
  p._element.pPr.set("algn",p_aliment)
  if paragraph_aliment=='left':
    p._element.pPr.set("indent",str(Pt(indent)))  #首行缩进，防止图文混排

  p.space_before,p.space_after,=paragraph.paragraph_format.space_before,paragraph.paragraph_format.space_after  #段落间距
  p.line_spacing=line_spacing

  # p.line_spacing=paragraph.paragraph_format.line_spacing
  #制表位
  if paragraph.paragraph_format.tab_stops:
    child_1 =  pptXML('a:tabLst')#制作一个标签
    for tab in paragraph.paragraph_format.tab_stops:
      child_2= pptXML('a:tab')
      child_2.set("pos",str(tab.position))
      child_1.append(child_2)
    p._element.pPr.append(child_1)
def set_table(newtable,table,rows_height,num,log):
  newtable._tbl.tblPr.remove(newtable._tbl.tblPr.xpath("a:tableStyleId")[0]) # 新表删除样式

  for newcolumn,column in zip(newtable.columns,table.columns):
      newcolumn.width=column.width   #设置新表列宽

  #通过pdf获取的
  for newrowtag,row_height in zip(newtable._tbl.xpath("a:tr"),rows_height
):
    newrowtag.set("h",str(Pt(row_height))) # 设置行高，通常为0

  # 单元格   跨几行，几列
  for i in range(len(table.rows)):
      for j in range(len(table.columns)):
        
        try:
          sorc=table._tbl.xpath(f"w:tr[{i+1}]/w:tc[{j+1}]/w:tcPr/w:vMerge/@w:val")[0]
          # table.cell(i,j)._tc.tcPr.xpath('w:vMerge')[0].get(wdqn("w:val"))
          if sorc=="continue":
            try:
              newtable.cell(i-1,j).merge(newtable.cell(i,j))
              continue
            except:
              log.write(f'第{num+1}页存在合并单元格，请手动调整\n')
        except IndexError:
          pass
        try:
          grid_span=table.cell(i,j).grid_span
          if not newtable.cell(i,j).is_spanned:
            if grid_span >1 :
              newtable.cell(i,j).merge(newtable.cell(i,j+grid_span-1))
              
            newcell=newtable.cell(i,j)
            newcell.margin_left=14604  # 左本框边距 常量
            newcell.margin_right=14604 # 右文本框边距
            newcell.margin_top=14604 # 上文本框边距
            newcell.margin_bottom=14604 # 下文本框边距
            # 填充 shape.Table.Cell(ii, jj).Borders(m).ForeColor.RGB = 0
            # 四种框线 宽度 线型 颜色
            cell_frames =["lnR","lnT","lnB","lnL"]
            for k in cell_frames:
              newcell._tc.tcPr.append( parse_xml(f'''<a:{k} {ptspace("a")} w="6350" cap="flat" cmpd="sng">
                <a:solidFill>
                  <a:srgbClr val="000008"/>
                </a:solidFill>
                <a:prstDash val="solid"/>
              </a:{k}>'''))
            # 单元格对齐方式
            if table.cell(i,j).vertical_alignment ==WD_ALIGN_VERTICAL.BOTTOM:
                newcell.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
            elif table.cell(i,j).vertical_alignment ==WD_ALIGN_VERTICAL.CENTER:
                newcell.vertical_anchor =MSO_VERTICAL_ANCHOR.MIDDLE
            elif table.cell(i,j).vertical_alignment ==WD_ALIGN_VERTICAL.TOP:
                newcell.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            text_frame=newcell.text_frame
            for a in range(len(table.cell(i,j).paragraphs)):
              if a==0:
                p =text_frame.paragraphs[0]
              else:
                p =text_frame.add_paragraph()
              for ele in table.cell(i,j).paragraphs[a]._element:
                  addelement(p,ele)
              set_paragraph(p,table.cell(i,j).paragraphs[a],0,1)
        except:
          log.write(f'第{num+1}页存在合并单元格，请手动调整\n')
  '''
  # 四种文本框边距   
  # 四种框线 宽度 线型 颜色
  # 填充
  for newrowtag,rowtag in zip(newtable._tbl.xpath("a:tr"),table._element.xpath("w:tr")):
    newrowtag.set("h",rowtag.trPr.trHeight.get(wdqn("w:val"))) # 设置行高，通常为0
    cell_attrs ={"top":[],"left":[],"bottom":[],"right":[]}
    for a in  cell_attrs:
      cell_attrs[a].append(rowtag.xpath(f"w:tblPrEx/w:tblBorders/w:{a}/@w:val")[0]) # 获取单元格边框属性
      cell_attrs[a].append(rowtag.xpath(f"w:tblPrEx/w:tblBorders/w:{a}/@w:color")[0]) # 获取单元格边框颜色
      cell_attrs[a].append( rowtag.xpath(f"w:tblPrEx/w:tblBorders/w:{a}/@w:sz")[0]) # 获取单元格边框大小
      cell_attrs[a].append( rowtag.xpath(f"w:tblPrEx/w:tblCellMar/w:{a}/@w:w")[0]) # 获取单元格边框宽度
    for celltag in rowtag.xpath("w:tc"):
  '''

class format_doc():
  def __init__(self,doc_data):
    self.doc = Document(doc_data)
  def set_titlestyle(self):
    try:
      style=self.doc.styles.add_style('题型', WD_STYLE_TYPE.PARAGRAPH)
    except:
       pass
    for paragraph in self.doc.paragraphs:
      paragraph.paragraph_format.line_spacing=1.2 #设置行距 
      if paragraph.text:
        # 设置标题,手动选取相似文本设置首个字符,题型前加&&
        if paragraph.runs[0].style.name.startswith("标题"):
          paragraph.style = f'Heading {paragraph.runs[0].style.name[3]}'
        elif paragraph.text.startswith('&&'):
          paragraph.style='题型'
        #设置制表位置
        tab_stops = paragraph.paragraph_format.tab_stops
        tab_stops.clear_all()
        select_count=re.findall(r"\t[B-D]\.",paragraph.text)
        if len(select_count)  >1:
          tab_stops.add_tab_stop(wdPt(160))  #常量
          tab_stops.add_tab_stop(wdPt(320))
          tab_stops.add_tab_stop(wdPt(480))
        else:
            tab_stops.add_tab_stop(wdPt(320))
  def set_font(self,replace_fonts={"方正黑体_GBK":"黑体","方正书宋_GBK":"宋体","方正楷体_GBK":"楷体","方正仿宋_GBK":"仿宋","方正魏碑_GBK":"黑体","方正准圆_GBK":"黑体","方正魏碑_GBK":"黑体"},anscolor='洋红'):

    #删除无用字符版面控制字符
    # for secttag in find_all_elements(self.doc._element, wdqn('w:sectPr')):#删除分节符
    #   secttag.getparent().remove(secttag)
    # for symboltag in find_all_elements(self.doc._element, wdqn('w:br')):
    #   if symboltag.get(wdqn('w:type')) in ['page','column']: #删除分栏符，分页符
    #     symboltag.getparent().remove(symboltag)     

    #替换字体名称
    for rFont in find_all_elements(self.doc._element, wdqn('w:rFonts')):    
      if rFont.get(wdqn('w:eastAsia')) in replace_fonts:
        rFont.set(wdqn('w:eastAsia'),replace_fonts[rFont.get(wdqn('w:eastAsia'))])
    #更改字体大小
    for fontsizetag in find_all_elements(self.doc._element, wdqn('w:sz')):
      fontsizetag.set(wdqn('w:val'),'48') #设置字体大小为24磅

    #替换颜色
    colors={'青色':'ABD8DA','洋红':'A46AA6'}
    for colortag in find_all_elements(self.doc._element, wdqn('w:color')): #更改字体的颜色
      if colortag.get(wdqn('w:val'))==colors[anscolor]: #答案的颜色
        colortag.set(wdqn('w:val'),'FF0000')
      elif colortag.get(wdqn('w:val')) in ['ABD8DA','A46AA6']:  #其他颜色
        colortag.set(wdqn('w:val'),'00B0F0')
    for utag in find_all_elements(self.doc._element, wdqn('w:u')): #更改下划线的颜色
      utag.set(wdqn('w:color'),'000000')

  def set_images(self,shape_scale=1.5): #处理图片
    for inlineshape in self.doc.inline_shapes:
        inlineshape.height=int(shape_scale*inlineshape.height) #放大图片
        inlineshape.width=int(shape_scale*inlineshape.width)
        # 删除指定的图，需要输入图片名，比较繁琐
        # del_pics=('image2.jpeg',) #常量
        # rid=inlineshape._inline.graphic.graphicData.pic.blipFill.blip.embed
        # if  self.doc._part.rels[rid]._target.partname.filename.startswith(del_pics):
        #     inlineshape._inline.getparent().remove(inlineshape._inline)
  #设置word页面尺寸和页边距
  def set_page(self):
    section = self.doc.sections[0]
    #根据ppt所容纳文本框尺寸确定
    # if mom_ppt:
    #   prs = Presentation("母版.pptx")
    #   section.page_height =prs.slide_height
    #   section.page_width =prs.slide_width
    #   section.top_margin = wdPt(83)
    #   section.bottom_margin = wdPt(64)
    #   section.left_margin = wdPt(28)
    #   section.right_margin = wdPt(28)
    #   prs.slide_masters[0].shapes[0].top.pt
    # else:
    section.page_height = wdPt(540) #常量
    section.page_width = wdPt(720)
    section.top_margin = wdPt(83)
    section.bottom_margin = wdPt(64)
    section.left_margin = wdPt(28)
    section.right_margin = wdPt(28)

  def set_page_break(self):
    last_para,previousstyle="",""
    for para in self.doc.paragraphs:
      para.paragraph_format.keep_together = True  # 所有段落分页格式设置成段中不分页
      para.paragraph_format.widow_control = True # 所有段落分页格式设置成孤行控制

      currentstyle=para.style.name
      if  "Heading" in previousstyle and currentstyle in ["Normal",'题型'] :
        para.paragraph_format.page_break_before=True #段前分页
      
      elif  previousstyle=="Normal" and ("Heading" in currentstyle or currentstyle=='题型'):
        para.paragraph_format.page_break_before=True #段前分页

      elif currentstyle=="题型" and previousstyle!="题型":
        para.paragraph_format.page_break_before=True
        para.paragraph_format.keep_with_next=True #与下段同
      previousstyle=currentstyle

      if currentstyle=="Normal":   #题目内容尽可能在一起
        if  re.match(r"\d{1,3}\.", para.text) or re.match(r"[一二三四五六七八]、", para.text): #题开始与下段分页
          para.paragraph_format.keep_with_next=True #与下段同页
          try:
            last_para.paragraph_format.keep_with_next=False
          except:
            pass
        last_para=para
  def mark_paragraph(self):# 添加^标记 ，为了分清楚是段落，图片还是表格
    for paragraph in self.doc.paragraphs:
      for run in paragraph._element:
        if  (run.tag==wdqn("w:r")   and run.text) or  run.tag==wdqn("m:oMath") or run.tag==wdqn("m:oMathPara"):
          paragraph._element.insert(paragraph._element.index(run), wdparse_xml(f'''<w:r {wdspace("w")}><w:rPr><w:rFonts w:hint="eastAsia"/><w:sz w:val="40"/><w:szCs w:val="40"/><w:lang w:val="en-US" w:eastAsia="zh-CN"/></w:rPr><w:t>^</w:t> </w:r>'''))
          break

  def mark_title(self):
    for paragraph in self.doc.paragraphs:
        if paragraph.style.name.startswith("Heading"):
          paragraph.text=f"&&{paragraph.text}"
  def omml2latex(self):
    oMaths= find_all_elements(self.doc._element, wdqn('m:oMath'))
    for oMath in oMaths:
      latex=omath_convert(oMath)
      runele=wdparse_xml(f'''<w:r {wdspace("w")}> <w:t>{latex}</w:t> </w:r>''')
      if oMath.getparent().tag==wdqn("m:oMathPara"):
        oMath.getparent().getparent().replace(oMath.getparent(),runele)
      else:
        oMath.getparent().replace(oMath,runele)
  def save(self):
    output = BytesIO()
    self.doc.save(output)
    return output.getvalue()

